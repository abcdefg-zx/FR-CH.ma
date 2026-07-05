from fastapi import FastAPI, HTTPException, UploadFile, File, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import requests
import os
import csv
import io
from datetime import datetime
from database import (
    get_memory_by_text,
    add_memory,
    get_all_memories,
    delete_memory,
    clear_all_memories,
    get_all_terminology,
    add_terminology,
    delete_terminology,
    clear_all_terminology,
    find_matching_terms,
    add_vocabulary,
    get_all_vocabulary,
    get_unmastered_vocabulary,
    update_vocabulary_mastered,
    delete_vocabulary,
    clear_all_vocabulary,
    update_vocabulary,
    add_vocabulary_meaning,
)

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

app = FastAPI(title="中法互译工具 API", version="2.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class TranslateRequest(BaseModel):
    text: str
    source_lang: str
    target_lang: str
    api_key: str = ""
    model: str = "deepseek-chat"
    base_url: str = ""

class ImportTermsRequest(BaseModel):
    terms: list[dict]

class TranscribeRequest(BaseModel):
    text: str
    source_lang: str
    target_lang: str
    api_key: str = ""
    model: str = "deepseek-chat"
    base_url: str = ""

class VocabularyUpdateRequest(BaseModel):
    mastered: int

def segment_text(text, source_lang):
    paragraphs = []
    
    if source_lang == 'zh':
        text = text.replace('\r\n', '\n')
        raw_paragraphs = text.split('\n\n')
        
        for para in raw_paragraphs:
            para = para.strip()
            if para:
                sentences = []
                current_sentence = ''
                
                for char in para:
                    current_sentence += char
                    if char in '。！？；\n':
                        sentences.append(current_sentence.strip())
                        current_sentence = ''
                
                if current_sentence.strip():
                    sentences.append(current_sentence.strip())
                
                paragraphs.append({
                    'text': para,
                    'sentences': sentences
                })
    else:
        text = text.replace('\r\n', '\n')
        raw_paragraphs = text.split('\n\n')
        
        for para in raw_paragraphs:
            para = para.strip()
            if para:
                sentences = []
                current_sentence = ''
                
                for char in para:
                    current_sentence += char
                    if char in '.!?;\n':
                        sentences.append(current_sentence.strip())
                        current_sentence = ''
                
                if current_sentence.strip():
                    sentences.append(current_sentence.strip())
                
                paragraphs.append({
                    'text': para,
                    'sentences': sentences
                })
    
    return paragraphs

def extract_vocabulary(source_text, target_text, source_lang, target_lang):
    if source_lang == 'zh':
        chinese_text = source_text
        french_text = target_text
    else:
        chinese_text = target_text
        french_text = source_text
    
    chinese_words = [word.strip() for word in chinese_text.replace('，', ' ').replace('。', ' ').replace('！', ' ').replace('？', ' ').split() if len(word.strip()) >= 2]
    french_words = []
    for word in french_text.replace(',', ' ').replace('.', ' ').replace('!', ' ').replace('?', ' ').split():
        clean_word = word.strip().lower()
        if len(clean_word) >= 3:
            french_words.append(clean_word)
    
    for i in range(min(len(chinese_words), len(french_words))):
        add_vocabulary(french_words[i], 'fr', chinese_words[i], 'zh', source_text[:100])

def translate_with_deepseek(text, source_lang, target_lang, api_key="", model="deepseek-chat", base_url="", terminology_hint=""):
    """DeepSeek 翻译 - Chat Completions API"""
    if source_lang == 'zh':
        source_name, target_name = "中文", "法语"
    else:
        source_name, target_name = "法语", "中文"
    
    system_content = f"你是专业的{source_name}到{target_name}翻译助手，请准确翻译，只输出译文。"
    if terminology_hint:
        system_content += f"\n请参考以下术语翻译：{terminology_hint}"
    
    url = base_url or "https://api.deepseek.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": model or "deepseek-chat",
        "messages": [
            {"role": "system", "content": system_content},
            {"role": "user", "content": text}
        ],
        "temperature": 0.3
    }
    
    response = requests.post(url, headers=headers, json=payload, timeout=60)
    response.raise_for_status()
    result = response.json()
    return result["choices"][0]["message"]["content"].strip()

def call_translation_api(text, source_lang, target_lang, api_key, model="deepseek-chat", base_url="", terminology_hint=""):
    """统一翻译入口 - DeepSeek 单引擎"""
    text = text.strip()
    if not text:
        return ""

    # 优先使用前端传入的 Key，否则回退到环境变量
    effective_key = api_key or os.getenv("DEEPSEEK_API_KEY", "")
    if not effective_key:
        raise HTTPException(status_code=400, detail="请先在「API设置」中填写 DeepSeek 的 API Key，或在服务器配置 DEEPSEEK_API_KEY 环境变量")
    
    try:
        return translate_with_deepseek(text, source_lang, target_lang, effective_key, model, base_url, terminology_hint)
    except HTTPException:
        raise
    except requests.exceptions.HTTPError as e:
        raise HTTPException(status_code=502, detail=f"DeepSeek API 调用失败: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"DeepSeek 翻译失败: {str(e)}")

def extract_text_from_docx(content):
    if not DOCX_AVAILABLE:
        raise Exception("python-docx not installed")
    
    doc = Document(io.BytesIO(content))
    paragraphs = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    
    return "\n\n".join(paragraphs)

def extract_text_from_pptx(content):
    if not PPTX_AVAILABLE:
        raise Exception("python-pptx not installed")
    
    prs = Presentation(io.BytesIO(content))
    paragraphs = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            paragraphs.append(f"[幻灯片 {slide_num}]\n" + "\n".join(slide_texts))
    
    return "\n\n".join(paragraphs)


@app.post("/api/translate")
async def translate(request: TranslateRequest):
    if not request.text.strip():
        raise HTTPException(status_code=400, detail="翻译文本不能为空")
    
    memory_result = get_memory_by_text(request.text.strip(), request.source_lang)
    if memory_result:
        target_text, target_lang, similarity = memory_result
        is_exact = similarity >= 1.0
        return {
            "success": True,
            "result": target_text,
            "from_memory": True,
            "is_exact_match": is_exact,
            "similarity": round(similarity, 2),
            "used_terminology": []
        }
    
    matching_terms = find_matching_terms(request.text, request.source_lang)
    terminology_hint = ", ".join(matching_terms) if matching_terms else ""
    
    translated_text = call_translation_api(
        request.text,
        request.source_lang,
        request.target_lang,
        request.api_key,
        request.model,
        request.base_url,
        terminology_hint
    )
    
    add_memory(request.text, request.source_lang, translated_text, request.target_lang)

    return {
        "success": True,
        "result": translated_text,
        "from_memory": False,
        "is_exact_match": False,
        "similarity": 0.0,
        "used_terminology": matching_terms
    }

@app.get("/api/memory")
async def get_memory():
    memories = get_all_memories()
    return {"success": True, "data": memories}

@app.delete("/api/memory/{memory_id}")
async def delete_memory_by_id(memory_id: int):
    success = delete_memory(memory_id)
    if not success:
        raise HTTPException(status_code=404, detail="记忆记录不存在")
    return {"success": True}

@app.delete("/api/memory")
async def clear_memory():
    clear_all_memories()
    return {"success": True}

@app.get("/api/terminology")
async def get_terminology():
    terms = get_all_terminology()
    return {"success": True, "data": terms}

@app.post("/api/terminology/import")
async def import_terminology(request: ImportTermsRequest):
    for term in request.terms:
        if "chinese_term" in term and "french_term" in term:
            add_terminology(term["chinese_term"], term["french_term"])
    return {"success": True, "count": len(request.terms)}

@app.post("/api/terminology/upload")
async def upload_terminology(file: UploadFile = File(...)):
    content = await file.read()
    rows = []
    
    if file.filename.endswith('.csv'):
        reader = csv.reader(io.StringIO(content.decode('utf-8')))
        for row in reader:
            if len(row) >= 2:
                rows.append([row[0].strip(), row[1].strip()])
    elif file.filename.endswith('.xlsx') or file.filename.endswith('.xls'):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True)
            ws = wb.active
            for row in ws.iter_rows(values_only=True):
                if len(row) >= 2:
                    rows.append([str(row[0]).strip() if row[0] else '', str(row[1]).strip() if row[1] else ''])
        except ImportError:
            raise HTTPException(status_code=400, detail="需要安装openpyxl库来处理Excel文件")
    else:
        raise HTTPException(status_code=400, detail="不支持的文件格式，请上传CSV或Excel文件")
    
    if not rows:
        raise HTTPException(status_code=400, detail="文件中没有有效的数据")
    
    count = 0
    for chinese, french in rows:
        if chinese and french and chinese != 'None' and french != 'None':
            add_terminology(chinese, french)
            count += 1
    
    return {"success": True, "count": count}

@app.delete("/api/terminology/{term_id}")
async def delete_term(term_id: int):
    success = delete_terminology(term_id)
    if not success:
        raise HTTPException(status_code=404, detail="术语不存在")
    return {"success": True}

@app.delete("/api/terminology")
async def clear_terminology():
    clear_all_terminology()
    return {"success": True}

@app.post("/api/transcribe")
async def transcribe(request: TranscribeRequest):
    if not request.text.strip():
        return {"success": False, "result": ""}
    
    memory_result = get_memory_by_text(request.text.strip(), request.source_lang)
    if memory_result:
        return {
            "success": True,
            "result": memory_result[0],
            "from_memory": True
        }
    
    matching_terms = find_matching_terms(request.text, request.source_lang)
    terminology_hint = ", ".join(matching_terms) if matching_terms else ""
    
    translated_text = call_translation_api(
        request.text,
        request.source_lang,
        request.target_lang,
        request.api_key,
        request.model,
        request.base_url,
        terminology_hint
    )
    
    add_memory(request.text, request.source_lang, translated_text, request.target_lang)
    
    return {
        "success": True,
        "result": translated_text,
        "from_memory": False
    }

@app.post("/api/document/translate")
async def translate_document(file: UploadFile = File(...), source_lang: str = "zh", target_lang: str = "fr", api_key: str = "", model: str = "deepseek-chat", base_url: str = ""):
    filename = file.filename.lower()
    content = await file.read()
    
    if filename.endswith('.txt'):
        text = content.decode('utf-8')
    elif filename.endswith('.csv'):
        reader = csv.reader(io.StringIO(content.decode('utf-8')))
        text = '\n'.join([','.join(row) for row in reader])
    elif filename.endswith('.xlsx') or filename.endswith('.xls'):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True)
            ws = wb.active
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(cell) if cell else '' for cell in row])
            text = '\n'.join(['\t'.join(row) for row in rows])
        except ImportError:
            raise HTTPException(status_code=400, detail="需要安装openpyxl库来处理Excel文件")
    elif filename.endswith('.docx') or filename.endswith('.doc'):
        try:
            text = extract_text_from_docx(content)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Word文档解析失败: {str(e)}")
    elif filename.endswith('.pptx') or filename.endswith('.ppt'):
        try:
            text = extract_text_from_pptx(content)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PPT文档解析失败: {str(e)}")
    else:
        raise HTTPException(status_code=400, detail="不支持的文件格式，请上传TXT、CSV、Excel、Word或PPT文件")
    
    paragraphs = segment_text(text, source_lang)
    
    translated_paragraphs = []
    total_sentences = 0
    translated_sentences = 0
    
    for para in paragraphs:
        translated_sentences_list = []
        
        for sentence in para['sentences']:
            total_sentences += 1
            
            memory_result = get_memory_by_text(sentence, source_lang)
            if memory_result:
                translated_sentences_list.append(memory_result[0])
                translated_sentences += 1
                continue
            
            matching_terms = find_matching_terms(sentence, source_lang)
            terminology_hint = ", ".join(matching_terms) if matching_terms else ""
            
            try:
                translated = call_translation_api(
                    sentence, source_lang, target_lang,
                    api_key, model, base_url, terminology_hint
                )
                
                add_memory(sentence, source_lang, translated, target_lang)
                translated_sentences_list.append(translated)
                translated_sentences += 1
            except Exception as e:
                translated_sentences_list.append(f"[翻译失败: {str(e)}]")
        
        translated_paragraphs.append({
            "original": para['text'],
            "translated": "\n".join(translated_sentences_list),
            "sentences_count": len(para['sentences'])
        })
    
    full_translation = "\n\n".join([p['translated'] for p in translated_paragraphs])
    
    return {
        "success": True,
        "filename": file.filename,
        "paragraphs": translated_paragraphs,
        "full_translation": full_translation,
        "stats": {
            "total_sentences": total_sentences,
            "translated_sentences": translated_sentences
        }
    }

@app.get("/api/export")
async def export_translations():
    memories = get_all_memories()
    
    export_data = []
    for m in memories:
        export_data.append({
            "id": m['id'],
            "source_text": m['source_text'],
            "target_text": m['target_text'],
            "source_lang": m['source_lang'],
            "target_lang": m['target_lang'],
            "created_at": m['created_at']
        })
    
    return {"success": True, "data": export_data}

@app.get("/api/export/txt")
async def export_translations_txt():
    memories = get_all_memories()
    
    content = "=== 中法翻译记录导出 ===\n\n"
    content += f"导出时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
    content += f"总记录数: {len(memories)}\n\n"
    content += "=" * 50 + "\n\n"
    
    for i, m in enumerate(memories, 1):
        direction = "中文→法语" if m['source_lang'] == 'zh' else "法语→中文"
        content += f"[{i}] {direction}\n"
        content += f"原文: {m['source_text']}\n"
        content += f"译文: {m['target_text']}\n"
        content += f"时间: {m['created_at']}\n"
        content += "-" * 30 + "\n\n"
    
    return Response(
        content=content.encode('utf-8'),
        media_type='text/plain',
        headers={
            'Content-Disposition': 'attachment; filename="translation_export.txt"'
        }
    )

class ExportSingleRequest(BaseModel):
    source_text: str
    target_text: str
    source_lang: str
    target_lang: str

@app.post("/api/export/single")
async def export_single_translation(request: ExportSingleRequest):
    direction = "中文→法语" if request.source_lang == 'zh' else "法语→中文"
    
    content = "=== 中法翻译结果导出 ===\n\n"
    content += f"导出时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
    content += f"翻译方向: {direction}\n\n"
    content += "=" * 50 + "\n\n"
    content += f"原文:\n{request.source_text}\n\n"
    content += "-" * 50 + "\n\n"
    content += f"译文:\n{request.target_text}\n"
    
    return Response(
        content=content.encode('utf-8'),
        media_type='text/plain',
        headers={
            'Content-Disposition': 'attachment; filename="current_translation.txt"'
        }
    )

@app.get("/api/vocabulary")
async def get_vocabulary(unmastered_only: bool = False):
    if unmastered_only:
        vocab = get_unmastered_vocabulary()
    else:
        vocab = get_all_vocabulary()
    return {"success": True, "data": vocab}

@app.put("/api/vocabulary/{vocab_id}")
async def update_vocab(vocab_id: int, request: VocabularyUpdateRequest):
    success = update_vocabulary_mastered(vocab_id, request.mastered)
    if not success:
        raise HTTPException(status_code=404, detail="生词不存在")
    return {"success": True}

@app.delete("/api/vocabulary/{vocab_id}")
async def delete_vocab(vocab_id: int):
    success = delete_vocabulary(vocab_id)
    if not success:
        raise HTTPException(status_code=404, detail="生词不存在")
    return {"success": True}

@app.delete("/api/vocabulary")
async def clear_vocabulary():
    clear_all_vocabulary()
    return {"success": True}

@app.post("/api/vocabulary/add")
async def add_vocab(request: dict):
    if "source_word" not in request or "target_word" not in request:
        raise HTTPException(status_code=400, detail="缺少必要参数")

    source_word = request["source_word"]
    source_lang = request.get("source_lang", "fr")
    target_word = request["target_word"]
    target_lang = request.get("target_lang", "zh")

    # 生成例句（使用模板，不依赖翻译引擎）
    example_sentence = ""
    if source_lang == 'fr':
        example_sentence = f"法语：{source_word}\n例句：{source_word} — {target_word}。\n中文：{target_word}"
    else:
        example_sentence = f"中文：{source_word}\n例句：{source_word} — {target_word}。\n法语：{target_word}"

    add_vocabulary(
        source_word, source_lang, target_word, target_lang,
        request.get("context", ""), example_sentence
    )
    return {"success": True, "example_sentence": example_sentence}


class GenerateExampleRequest(BaseModel):
    source_word: str
    source_lang: str
    target_word: str

@app.post("/api/vocabulary/generate-example")
async def generate_example(request: GenerateExampleRequest):
    # 使用模板生成例句，不依赖翻译引擎
    if request.source_lang == 'fr':
        example = f"法语：{request.source_word}\n例句：{request.source_word} — {request.target_word}。\n中文：{request.target_word}"
    else:
        example = f"中文：{request.source_word}\n例句：{request.source_word} — {request.target_word}。\n法语：{request.target_word}"
    return {"success": True, "example": example}


class EditVocabRequest(BaseModel):
    source_word: str = None
    target_word: str = None
    example_sentence: str = None

@app.put("/api/vocabulary/{vocab_id}/edit")
async def edit_vocabulary(vocab_id: int, request: EditVocabRequest):
    success = update_vocabulary(vocab_id, request.source_word, request.target_word, request.example_sentence)
    if not success:
        raise HTTPException(status_code=404, detail="生词不存在")
    return {"success": True}


class AddMeaningRequest(BaseModel):
    new_meaning: str

@app.post("/api/vocabulary/{vocab_id}/meaning")
async def add_meaning(vocab_id: int, request: AddMeaningRequest):
    success = add_vocabulary_meaning(vocab_id, request.new_meaning)
    if not success:
        raise HTTPException(status_code=404, detail="生词不存在")
    return {"success": True}


class ValidateTranslationRequest(BaseModel):
    source_text: str
    user_translation: str
    source_lang: str
    target_lang: str

@app.post("/api/translation/validate")
async def validate_translation(request: ValidateTranslationRequest):
    # 无免费引擎可用，直接保留用户修改
    return {"success": True, "valid": True, "message": "已保留您的修改", "engine_translation": ""}


class ExportDocumentRequest(BaseModel):
    filename: str = "document"
    paragraphs: list[dict]

@app.post("/api/document/export")
async def export_document(request: ExportDocumentRequest):
    content = f"=== 文档翻译对照导出 ===\n\n"
    content += f"导出时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
    content += f"段落数: {len(request.paragraphs)}\n\n"
    content += "=" * 60 + "\n\n"

    for i, p in enumerate(request.paragraphs, 1):
        content += f"【段落 {i}】\n"
        content += f"原文：\n{p.get('original', '')}\n\n"
        content += f"译文：\n{p.get('translated', '')}\n\n"
        content += "-" * 60 + "\n\n"

    return Response(
        content=content.encode('utf-8'),
        media_type='text/plain',
        headers={
            'Content-Disposition': f'attachment; filename="document_translation.txt"'
        }
    )

frontend_dir = os.path.join(os.path.dirname(__file__), "../frontend")
if os.path.exists(frontend_dir):
    app.mount("/", StaticFiles(directory=frontend_dir, html=True), name="frontend")
else:
    @app.get("/")
    async def root():
        return {"message": "中法互译工具 API", "version": "2.0.0"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)